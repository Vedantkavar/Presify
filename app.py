import streamlit as st
# Set page config must be the very first Streamlit command
st.set_page_config(page_title="PDF to PowerPoint Converter", layout="wide")

# Fix SSL certificate issues
import ssl
try:
    _create_unverified_https_context = ssl._create_unverified_context
except AttributeError:
    pass
else:
    ssl._create_default_https_context = _create_unverified_https_context

import PyPDF2
import io
import re
import nltk
import sys
import traceback
from nltk.corpus import stopwords
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.probability import FreqDist
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
import base64
from collections import Counter
import random
import textwrap
from sklearn.feature_extraction.text import TfidfVectorizer

# Download necessary NLTK resources with better error handling
@st.cache_resource
def download_nltk_resources():
    try:
        # Check if data already exists first
        try:
            nltk.data.find('tokenizers/punkt')
            nltk.data.find('corpora/stopwords')
            st.success("NLTK resources already available")
        except LookupError:
            # If not found, download with SSL verification disabled
            st.info("Downloading NLTK resources...")
            nltk.download('punkt', quiet=True)
            nltk.download('stopwords', quiet=True)
            st.success("NLTK resources downloaded successfully")
    except Exception as e:
        st.warning(f"NLTK resource download issue: {e}")
        st.info("Will use fallback methods that don't require NLTK resources")

# Initialize fallback variables
FALLBACK_MODE = False

# Try to download NLTK data
try:
    download_nltk_resources()
except Exception as e:
    st.warning(f"Switching to fallback mode due to NLTK error: {e}")
    FALLBACK_MODE = True

def extract_text_from_pdf(pdf_file):
    try:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text += page.extract_text()
        return text
    except Exception as e:
        st.error(f"Error extracting text from PDF: {e}")
        return ""

def preprocess_text(text):
    # Clean the text
    text = re.sub(r'\s+', ' ', text)  # Remove extra whitespace
    text = re.sub(r'\n+', ' ', text)  # Remove newlines
    return text

def extract_title(text, max_length=60):
    try:
        if FALLBACK_MODE:
            # Simple fallback method
            first_line = text.split('.')[0]
            if len(first_line) > max_length:
                return first_line[:max_length] + "..."
            return first_line
        
        # Get first few sentences
        sentences = sent_tokenize(text)
        if not sentences:
            return "PDF Summary"
        
        # Try to find a title in the first few sentences
        potential_titles = []
        for i in range(min(3, len(sentences))):
            # Look for short, capitalized phrases that might be titles
            if len(sentences[i]) < max_length and any(w.isupper() for w in sentences[i].split()):
                potential_titles.append(sentences[i])
        
        # If we found potential titles, use the first one
        if potential_titles:
            return potential_titles[0]
        
        # Otherwise, use the first sentence or part of it
        first_sentence = sentences[0]
        if len(first_sentence) > max_length:
            return first_sentence[:max_length] + "..."
        return first_sentence
    except Exception as e:
        st.warning(f"Title extraction error: {e}. Using default title.")
        return "PDF Summary"

def identify_sections(text):
    """Identify potential sections in the document based on formatting and content"""
    try:
        # Look for common section indicators
        section_patterns = [
            r'(?i)chapter\s+\d+[:\s]+(.*?)(?=\n|$)',
            r'(?i)section\s+\d+[:\s]+(.*?)(?=\n|$)',
            r'(?i)^\s*\d+\.\s+(.*?)(?=\n|$)',
            r'(?i)^\s*[IVX]+\.\s+(.*?)(?=\n|$)',
            r'(?i)^\s*[A-Z][A-Z\s]+[A-Z](?:\s*\n|\s*$|\s*:)',  # ALL CAPS HEADINGS
            r'(?i)introduction|conclusion|background|methodology|results|discussion|summary|references'
        ]
        
        potential_sections = []
        for pattern in section_patterns:
            matches = re.finditer(pattern, text, re.MULTILINE)
            for match in matches:
                if match.group(0):
                    clean_heading = re.sub(r'^\s*\d+\.\s*', '', match.group(0))
                    clean_heading = re.sub(r'[:\n]', '', clean_heading).strip()
                    if 3 < len(clean_heading) < 60:  # Reasonable length for a heading
                        potential_sections.append(clean_heading)
        
        # If we don't find enough sections, create general ones
        if len(potential_sections) < 3:
            # Create generic section titles
            generic_titles = [
                "Title Slide",
                "Introduction & Background",
                "Research Objectives",
                "Literature Review & Theoretical Framework",
                "Research Methodology",
                "Key Findings & Results",
                "Data Analysis & Insights",
                "Discussion & Implications",
                "Conclusions & Recommendations",
                "Future Directions",
                "References & Acknowledgments",
                "Q&A"
            ]

            
            if FALLBACK_MODE:
                return generic_titles[:5]
            
            sentences = sent_tokenize(text)
            chunk_size = max(1, len(sentences) // 5)  # Divide into about 5 sections
            potential_sections = generic_titles[:min(5, len(sentences) // chunk_size + 1)]
        
        return potential_sections[:8]  # Limit to 8 sections for a reasonable presentation
    except Exception as e:
        st.warning(f"Section identification error: {e}. Using default sections.")
        return [
            "Introduction & Background",
            "Key Concepts",
            "Main Findings",
            "Conclusions & Recommendations"
        ]

def extract_key_points(text, section_title, num_points=4):
    """Extract key points related to a section title using TF-IDF or fallback to simple extraction"""
    try:
        if FALLBACK_MODE:
            # Simple fallback method without NLTK
            sentences = re.split(r'(?<=[.!?])\s+', text)
            return [s for s in sentences[:num_points * 2] if len(s) > 20][:num_points]
        
        sentences = sent_tokenize(text)
        
        # If very few sentences, just return them
        if len(sentences) <= num_points:
            return [s for s in sentences if len(s) > 20]  # Only return substantial sentences
        
        # Use TF-IDF to find relevant sentences
        try:
            # Create a vectorizer
            vectorizer = TfidfVectorizer(stop_words='english')
            tfidf_matrix = vectorizer.fit_transform(sentences)
            
            # Get feature names (words)
            feature_names = vectorizer.get_feature_names_out()
            
            # Calculate scores for each sentence
            scores = []
            query_words = word_tokenize(section_title.lower())
            stop_words = set(stopwords.words('english'))
            query_words = [w for w in query_words if w.isalnum() and w not in stop_words]
            
            for i, sentence in enumerate(sentences):
                score = 0
                # Get the TF-IDF scores for this sentence
                feature_index = tfidf_matrix[i, :].nonzero()[1]
                tfidf_scores = zip(feature_index, [tfidf_matrix[i, x] for x in feature_index])
                
                # Check if any query words are in the top TF-IDF terms
                for idx, score_value in tfidf_scores:
                    term = feature_names[idx]
                    if term in query_words:
                        score += score_value
                
                # Also consider sentence length and position
                if 20 < len(sentence) < 200:  # Prefer medium-length sentences
                    score += 0.1
                
                scores.append((i, score, sentence))
            
            # Get top sentences
            top_sentences = sorted(scores, key=lambda x: x[1], reverse=True)[:num_points]
            top_sentences = sorted(top_sentences, key=lambda x: x[0])  # Sort by original order
            
            key_points = [s[2] for s in top_sentences if len(s[2]) > 20]
            
            # If we didn't get enough points, add some more sentences
            if len(key_points) < num_points:
                additional_sentences = [s for s in sentences if s not in key_points and len(s) > 20]
                key_points.extend(additional_sentences[:num_points - len(key_points)])
            
            return key_points[:num_points]
        
        except Exception as e:
            # Fallback method if TF-IDF fails
            st.warning(f"Using simple extraction method due to: {e}")
            return [s for s in sentences[:num_points * 2] if len(s) > 20][:num_points]
    
    except Exception as e:
        st.warning(f"Key point extraction error: {e}. Using simple method.")
        # Very simple fallback
        sentences = re.split(r'(?<=[.!?])\s+', text)
        return [s for s in sentences[:num_points * 2] if len(s) > 20][:num_points]

def create_presentation(pdf_title, sections, section_contents):
    """Create PowerPoint presentation with OWASP-inspired design"""
    try:
        prs = Presentation()
        
        # Define slide layouts
        title_slide_layout = prs.slide_layouts[0]  # Title slide
        section_title_layout = prs.slide_layouts[2]  # Section title
        content_layout = prs.slide_layouts[1]  # Title and content
        
        # OWASP-inspired color scheme
        owasp_colors = {
            'background': RGBColor(18, 14, 64),  # Dark navy/purple background
            'title': RGBColor(125, 230, 188),  # Mint green for titles
            'subtitle': RGBColor(255, 255, 255),  # White for subtitles
            'heading': RGBColor(125, 230, 188),  # Mint green for headings
            'accent': RGBColor(252, 129, 110),  # Coral for accents
            'text': RGBColor(255, 255, 255),  # White for text
            'bullet': RGBColor(252, 129, 110)  # Coral for bullets
        }
        
        # Add title slide
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Change background color to OWASP dark blue
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = owasp_colors['background']
        
        # Add dots pattern effect placeholder (this would be implemented via shapes in a real scenario)
        # Here we would add semi-transparent circles on left and right for the dot pattern
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        # Set title and subtitle
        title_shape.text = pdf_title.upper()
        subtitle_shape.text = "PDF SUMMARY\nPowered by AI"
        
        # Format title slide
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.color.rgb = owasp_colors['title']
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = owasp_colors['subtitle']
        
        # Add tagline similar to "Protecting Our Digital World"
        tf = subtitle_shape.text_frame
        p = tf.add_paragraph()
        p.text = "\nExtracting Knowledge From Documents"
        p.font.size = Pt(16)
        p.font.color.rgb = owasp_colors['subtitle']
        
        # Add a table of contents slide
        toc_slide = prs.slides.add_slide(content_layout)
        
        # Change background color 
        toc_slide.background.fill.solid()
        toc_slide.background.fill.fore_color.rgb = owasp_colors['background']
        
        toc_title = toc_slide.shapes.title
        toc_title.text = "CONTENTS"
        toc_title.text_frame.paragraphs[0].font.color.rgb = owasp_colors['title']
        toc_title.text_frame.paragraphs[0].font.bold = True
        
        toc_content = toc_slide.placeholders[1]
        tf = toc_content.text_frame
        
        for i, section in enumerate(sections):
            p = tf.add_paragraph()
            p.text = f"{i+1}. {section}"
            p.font.size = Pt(20)
            p.font.color.rgb = owasp_colors['text']
            p.space_after = Pt(15)
        
        # Add content slides for each section
        for i, (section, content) in enumerate(zip(sections, section_contents)):
            # Add section title slide
            section_slide = prs.slides.add_slide(section_title_layout)
            
            # Change background color
            section_slide.background.fill.solid()
            section_slide.background.fill.fore_color.rgb = owasp_colors['background']
            
            section_title = section_slide.shapes.title
            section_title.text = section.upper()
            section_title.text_frame.paragraphs[0].font.size = Pt(40)
            section_title.text_frame.paragraphs[0].font.color.rgb = owasp_colors['title']
            section_title.text_frame.paragraphs[0].font.bold = True
            
            # Add slide with bullet points
            bullet_slide = prs.slides.add_slide(content_layout)
            
            # Change background color
            bullet_slide.background.fill.solid()
            bullet_slide.background.fill.fore_color.rgb = owasp_colors['background']
            
            bullet_title = bullet_slide.shapes.title
            bullet_title.text = section.upper()
            bullet_title.text_frame.paragraphs[0].font.color.rgb = owasp_colors['title']
            bullet_title.text_frame.paragraphs[0].font.bold = True
            
            bullet_content = bullet_slide.placeholders[1]
            tf = bullet_content.text_frame
            
            # Add key points with accent-colored bullet markers
            for point in content:
                # Wrap long content
                wrapped_lines = textwrap.wrap(point, width=90)
                point_text = wrapped_lines[0]
                if len(wrapped_lines) > 1:
                    point_text += "..."
                
                p = tf.add_paragraph()
                p.text = point_text
                p.font.size = Pt(18)
                p.font.color.rgb = owasp_colors['text']
                p.level = 0
                
                # Can't directly change bullet color in python-pptx, would require COM interaction
        
        # Add a conclusion slide
        conclusion_slide = prs.slides.add_slide(content_layout)
        
        # Change background color
        conclusion_slide.background.fill.solid()
        conclusion_slide.background.fill.fore_color.rgb = owasp_colors['background']
        
        conclusion_title = conclusion_slide.shapes.title
        conclusion_title.text = "KEY TAKEAWAYS"
        conclusion_title.text_frame.paragraphs[0].font.color.rgb = owasp_colors['title']
        conclusion_title.text_frame.paragraphs[0].font.bold = True
        
        conclusion_content = conclusion_slide.placeholders[1]
        tf = conclusion_content.text_frame
        
        # Add a few general conclusion points
        conclusion_points = [
            "This presentation provided an overview of the key points extracted from the document",
            "Sections were organized based on the main topics identified in the content",
            "For more detailed information, please refer to the original PDF document"
        ]
        
        for point in conclusion_points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.font.color.rgb = owasp_colors['text']
            p.level = 0
        
        return prs
    except Exception as e:
        error_msg = f"Error in create_presentation: {str(e)}\n{traceback.format_exc()}"
        st.error(error_msg)
        raise Exception(error_msg)

def get_download_link(pptx_content, filename):
    try:
        b64 = base64.b64encode(pptx_content).decode()
        return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}" download="{filename}" class="download-button">Download PowerPoint Presentation</a>'
    except Exception as e:
        st.error(f"Error creating download link: {e}")
        return None

# Custom CSS for better appearance - OWASP-inspired styling
st.markdown("""
<style>
.main {
  padding: 2.5rem;
  background: linear-gradient(135deg, #0f0c29, #302b63, #24243e);
  min-height: 100vh;
}

.stButton button {
  background: linear-gradient(135deg, #12103f, #1a1758);
  color: #7de6bc;
  padding: 14px 28px;
  border-radius: 12px;
  border: none;
  font-size: 18px;
  font-weight: 600;
  letter-spacing: 0.5px;
  transition: all 0.4s cubic-bezier(0.175, 0.885, 0.32, 1.275);
  box-shadow: 0 10px 20px rgba(0, 0, 0, 0.15), 0 3px 6px rgba(0, 0, 0, 0.1);
  position: relative;
  overflow: hidden;
}

.stButton button:hover {
  background: linear-gradient(135deg, #1a1758, #12103f);
  box-shadow: 0 15px 25px rgba(0, 0, 0, 0.25), 0 5px 10px rgba(0, 0, 0, 0.15);
  transform: translateY(-3px);
}

.stButton button:active {
  transform: translateY(1px);
  box-shadow: 0 5px 10px rgba(0, 0, 0, 0.2);
}

.stButton button::after {
  content: "";
  position: absolute;
  top: 0;
  left: 0;
  width: 100%;
  height: 100%;
  background: linear-gradient(90deg, transparent, rgba(255, 255, 255, 0.1), transparent);
  transform: translateX(-100%);
  transition: 0.5s;
}

.stButton button:hover::after {
  transform: translateX(100%);
}

h1 {
  color: rgb(255, 255, 255);
  margin-bottom: 2.5rem;
  font-size: 3rem;
  font-weight: 800;
  letter-spacing: -0.5px;
  text-shadow: 0 2px 10px rgba(0, 0, 0, 0.2);
  background: linear-gradient(to right, #ffffff, #7de6bc);
  -webkit-background-clip: text;
  background-clip: text;
  color: transparent;
  animation: fadeIn 1s ease-out;
}

@keyframes fadeIn {
  from {
    opacity: 0;
    transform: translateY(-10px);
  }
  to {
    opacity: 1;
    transform: translateY(0);
  }
}

.download-button {
  display: inline-block;
  background: linear-gradient(135deg, #12103f, #1a1758);
  color: #7de6bc !important;
  padding: 16px 32px;
  text-align: center;
  text-decoration: none;
  font-size: 18px;
  border-radius: 12px;
  font-weight: 600;
  margin-top: 24px;
  transition: all 0.4s cubic-bezier(0.175, 0.885, 0.32, 1.275);
  box-shadow: 0 10px 20px rgba(0, 0, 0, 0.15), 0 3px 6px rgba(0, 0, 0, 0.1);
  position: relative;
  overflow: hidden;
  z-index: 1;
}

.download-button:hover {
  background: linear-gradient(135deg, #1a1758, #12103f);
  box-shadow: 0 15px 25px rgba(0, 0, 0, 0.25), 0 5px 10px rgba(0, 0, 0, 0.15);
  transform: translateY(-3px);
}

.download-button:active {
  transform: translateY(1px);
  box-shadow: 0 5px 10px rgba(0, 0, 0, 0.2);
}

.download-button::before {
  content: "";
  position: absolute;
  top: 0;
  left: 0;
  width: 100%;
  height: 100%;
  background: linear-gradient(45deg, rgba(125, 230, 188, 0.1), rgba(125, 230, 188, 0));
  z-index: -1;
  transform: scaleX(0);
  transform-origin: 0 50%;
  transition: transform 0.5s ease-out;
}

.download-button:hover::before {
  transform: scaleX(1);
}

.stSuccess {
  background: rgba(18, 16, 63, 0.05);
  backdrop-filter: blur(10px);
  padding: 20px;
  border-radius: 16px;
  border-left: 5px solid #7de6bc;
  box-shadow: 0 8px 32px rgba(0, 0, 0, 0.1);
  margin: 20px 0;
  transition: all 0.3s ease;
  border-top: 1px solid rgba(255, 255, 255, 0.1);
  border-right: 1px solid rgba(255, 255, 255, 0.05);
  animation: slideIn 0.5s ease-out;
}

.stWarning {
  background: rgba(18, 16, 63, 0.05);
  backdrop-filter: blur(10px);
  padding: 20px;
  border-radius: 16px;
  border-left: 5px solid #fc816e;
  box-shadow: 0 8px 32px rgba(0, 0, 0, 0.1);
  margin: 20px 0;
  transition: all 0.3s ease;
  border-top: 1px solid rgba(255, 255, 255, 0.1);
  border-right: 1px solid rgba(255, 255, 255, 0.05);
  animation: slideIn 0.5s ease-out;
}

.stError {
  background: rgba(18, 16, 63, 0.05);
  backdrop-filter: blur(10px);
  padding: 20px;
  border-radius: 16px;
  border-left: 5px solid #ff5555;
  box-shadow: 0 8px 32px rgba(0, 0, 0, 0.1);
  margin: 20px 0;
  transition: all 0.3s ease;
  border-top: 1px solid rgba(255, 255, 255, 0.1);
  border-right: 1px solid rgba(255, 255, 255, 0.05);
  animation: slideIn 0.5s ease-out;
}

@keyframes slideIn {
  from {
    opacity: 0;
    transform: translateX(-10px);
  }
  to {
    opacity: 1;
    transform: translateX(0);
  }
}

.debug-info {
  background: rgba(248, 249, 250, 0.05);
  backdrop-filter: blur(10px);
  padding: 16px;
  border-radius: 12px;
  margin-top: 24px;
  font-family: "JetBrains Mono", "Fira Code", monospace;
  box-shadow: 0 8px 32px rgba(0, 0, 0, 0.1);
  border: 1px solid rgba(255, 255, 255, 0.05);
  color: #e0e0e0;
  transition: all 0.3s ease;
}

.debug-info:hover {
  box-shadow: 0 12px 40px rgba(0, 0, 0, 0.15);
  transform: translateY(-2px);
}

/* Additional aesthetic enhancements */
::selection {
  background: rgba(125, 230, 188, 0.3);
  color: white;
}

body {
  font-family: "Inter", "Segoe UI", system-ui, -apple-system, sans-serif;
  color: #f0f0f0;
  line-height: 1.6;
}

/* Scrollbar styling */
::-webkit-scrollbar {
  width: 10px;
}

::-webkit-scrollbar-track {
  background: rgba(18, 16, 63, 0.2);
}

::-webkit-scrollbar-thumb {
  background: rgba(125, 230, 188, 0.5);
  border-radius: 10px;
}

::-webkit-scrollbar-thumb:hover {
  background: rgba(125, 230, 188, 0.7);
}


</style>
""", unsafe_allow_html=True)

def main():
    st.title("Enhanced PDF to PowerPoint Converter")
    st.markdown("Upload a PDF file to generate a professional PowerPoint presentation with intelligent content extraction and OWASP-inspired design.")

    
    uploaded_file = st.file_uploader("Choose a PDF file", type="pdf")
    
    if uploaded_file is not None:
        st.success("File uploaded successfully!")
        
        # Display PDF information
        file_details = {
            "Filename": uploaded_file.name,
            "File size": f"{uploaded_file.size / 1024:.2f} KB"
        }
        st.write("📄 File Details:")
        for key, value in file_details.items():
            st.write(f"- {key}: {value}")
        
        # Generate PPT button
        if st.button("Generate Professional PowerPoint"):
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            try:
                # Process the PDF
                status_text.text("Step 1/5: Extracting text from PDF...")
                progress_bar.progress(10)
                
                pdf_text = extract_text_from_pdf(uploaded_file)
                
                if not pdf_text.strip():
                    st.error("Could not extract text from the PDF. The file may be scanned or protected.")
                    return
                
                # Preprocess the text
                status_text.text("Step 2/5: Preprocessing text...")
                progress_bar.progress(30)
                processed_text = preprocess_text(pdf_text)
                
                # Extract title
                status_text.text("Step 3/5: Extracting document title...")
                progress_bar.progress(40)
                pdf_title = extract_title(processed_text)
                
                # Identify sections
                status_text.text("Step 4/5: Identifying document sections...")
                progress_bar.progress(60)
                sections = identify_sections(processed_text)
                
                # Extract key points for each section
                status_text.text("Step 4/5: Extracting key points...")
                progress_bar.progress(70)
                
                section_contents = []
                for section in sections:
                    points = extract_key_points(processed_text, section)
                    section_contents.append(points)
                
                # Create PowerPoint
                status_text.text("Step 5/5: Creating PowerPoint presentation...")
                progress_bar.progress(85)
                
                prs = create_presentation(pdf_title, sections, section_contents)
                
                # Save the presentation to a BytesIO object
                pptx_io = io.BytesIO()
                try:
                    prs.save(pptx_io)
                    pptx_io.seek(0)
                    progress_bar.progress(95)
                    
                    # Create download link
                    download_link = get_download_link(
                        pptx_io.getvalue(), 
                        f"{uploaded_file.name.split('.')[0]}_presentation.pptx"
                    )
                    
                    if download_link:
                        st.markdown(download_link, unsafe_allow_html=True)
                        progress_bar.progress(100)
                        status_text.text("Complete!")
                        st.success("✨ PowerPoint presentation generated successfully! Click the button below to download.")
                    else:
                        st.error("Failed to create download link.")
                        
                except Exception as e:
                    st.error(f"Error saving presentation: {str(e)}")
                
                # Show preview info
                with st.expander("📋 Presentation Content Preview"):
                    st.markdown(f"### Title: {pdf_title}")
                    st.markdown("### Sections:")
                    for i, (section, points) in enumerate(zip(sections, section_contents)):
                        st.markdown(f"**{i+1}. {section}**")
                        for j, point in enumerate(points):
                            st.markdown(f"   • {point[:100]}..." if len(point) > 100 else f"   • {point}")
            
            except Exception as e:
                st.error(f"An error occurred: {str(e)}")
                st.info("Please try with a different PDF file or check if the file is properly formatted.")

if __name__ == "__main__":
    main()